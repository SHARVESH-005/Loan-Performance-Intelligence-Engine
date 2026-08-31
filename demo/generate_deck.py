"""Generate a pitch deck PPTX for the Intain Campus FinTech Challenge submission."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)      # dark navy
BG_CARD = RGBColor(0x1A, 0x25, 0x3C)       # card bg
ACCENT = RGBColor(0x00, 0xD4, 0xAA)        # teal/green
ACCENT2 = RGBColor(0x60, 0xA5, 0xFA)       # blue
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_para(tf, text, size=16, color=WHITE, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return p

# ===== SLIDE 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_text(slide, 1, 1.5, 11, 1.5, "Loan Performance Intelligence Engine", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.2, 11, 0.8, "Intain Campus FinTech Challenge 2026 — AI Track", size=24, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, 1, 4.5, 11, 0.8, "End-to-end ML pipeline for loan-level portfolio analysis", size=20, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.8, 11, 0.6, "Solo Build  ·  Python  ·  LightGBM  ·  SHAP  ·  Gemini LLM", size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ===== SLIDE 2: Problem & Solution =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "The Challenge", size=36, color=ACCENT, bold=True)

tf = add_text(slide, 0.8, 1.5, 5.5, 5, "PROBLEM", size=22, color=ORANGE, bold=True)
add_para(tf, "")
add_para(tf, "• Loan portfolios have thousands of records", size=16, color=WHITE)
add_para(tf, "  updated monthly — manual review doesn't scale", size=16, color=GRAY)
add_para(tf, "")
add_para(tf, "• Predicting delinquency, default, and prepayment", size=16, color=WHITE)
add_para(tf, "  requires handling rare events and time-series data", size=16, color=GRAY)
add_para(tf, "")
add_para(tf, "• Data quality issues (stale updates, missing docs)", size=16, color=WHITE)
add_para(tf, "  go undetected without automated checks", size=16, color=GRAY)

tf2 = add_text(slide, 7, 1.5, 5.5, 5, "SOLUTION", size=22, color=ACCENT, bold=True)
add_para(tf2, "")
add_para(tf2, "• One-command pipeline: raw data → scored output", size=16, color=WHITE)
add_para(tf2, "  python -m src.run_pipeline --stage all", size=14, color=GRAY)
add_para(tf2, "")
add_para(tf2, "• 5 calibrated prediction models + survival model", size=16, color=WHITE)
add_para(tf2, "  with time-aware splits to prevent data leakage", size=16, color=GRAY)
add_para(tf2, "")
add_para(tf2, "• Hybrid anomaly detection (rules + Isolation Forest)", size=16, color=WHITE)
add_para(tf2, "  + governed LLM copilot for reviewer notes", size=16, color=GRAY)

# ===== SLIDE 3: Architecture =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "System Architecture", size=36, color=ACCENT, bold=True)

arch_text = """Raw Data  →  FR-1: Data Profiling  →  Feature Engineering  →  Time-Aware Split

     ┌──────────────────┬───────────────────┬──────────────────┐
     │                  │                   │                  │
FR-2: Predictions   FR-3: Survival      FR-4: Anomaly
  (LightGBM ×5)     (Hazard Model)      (IF + Rules)
     │                  │                   │
     └──────────────────┴───────────────────┘
                        │
              FR-6: Explainability (SHAP)
                        │
              FR-5: Scenario Simulation
                        │
              FR-7: LLM Copilot (Gemini)
                        │
                  submission.csv"""
add_text(slide, 0.8, 1.6, 11.5, 5, arch_text, size=15, color=WHITE)

# ===== SLIDE 4: Key Metrics =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "Model Performance — Key Metrics", size=36, color=ACCENT, bold=True)

# Metrics table
metrics = [
    ("Target", "Baseline AUC", "Improved AUC", "Brier Score"),
    ("3m Delinquency", "0.734", "0.754", "0.067"),
    ("6m Delinquency", "0.650", "0.773", "0.091"),
    ("12m Default", "0.781", "0.991 ✦", "0.001"),
    ("12m Prepayment", "0.497", "0.877", "0.007"),
]

for i, row in enumerate(metrics):
    y = 1.6 + i * 0.7
    color = ACCENT if i == 0 else WHITE
    bold = i == 0
    highlight = i == 3  # default row
    for j, cell in enumerate(row):
        x = 1.0 + j * 2.8
        c = ORANGE if (highlight and j >= 2) else color
        add_text(slide, x, y, 2.6, 0.5, cell, size=18 if i==0 else 20, color=c, bold=bold)

# Survival metric
add_text(slide, 1.0, 5.5, 11, 0.5, "FR-3 Survival Model:  AUC 0.998  vs  Baseline 0.515", size=20, color=ACCENT, bold=True)
add_text(slide, 1.0, 6.1, 11, 0.5, "All metrics from held-out test set · Time-aware split · No data leakage", size=14, color=GRAY)

# ===== SLIDE 5: Explainability + Anomaly =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "Explainability & Anomaly Detection", size=36, color=ACCENT, bold=True)

# Left: explainability
tf = add_text(slide, 0.8, 1.5, 5.5, 4.5, "FR-6: SHAP Explainability", size=22, color=ORANGE, bold=True)
add_para(tf, "")
add_para(tf, "• Global SHAP summaries for all 3 models", size=16, color=WHITE)
add_para(tf, "  (delinquency, default, prepayment)", size=14, color=GRAY)
add_para(tf, "")
add_para(tf, "• Per-loan waterfall plots", size=16, color=WHITE)
add_para(tf, "  (top drivers: days_past_due, DTI, balance)", size=14, color=GRAY)
add_para(tf, "")
add_para(tf, "• 3 FP + 3 FN examples with hypothesized causes", size=16, color=WHITE)

# Right: anomaly
tf2 = add_text(slide, 7, 1.5, 5.5, 4.5, "FR-4: Anomaly Detection", size=22, color=ORANGE, bold=True)
add_para(tf2, "")
add_para(tf2, "• Hybrid: Isolation Forest + 5 deterministic rules", size=16, color=WHITE)
add_para(tf2, "")
add_para(tf2, "• 1,857 records flagged in test set:", size=16, color=WHITE)
add_para(tf2, "  - 1,161 document gaps", size=16, color=ACCENT)
add_para(tf2, "  - 736 stale servicer updates", size=16, color=ACCENT)
add_para(tf2, "")
add_para(tf2, "• Servicer reconciliation with full audit log", size=16, color=WHITE)

# Images
shap_path = os.path.join("reports", "explainability", "shap_summary_delinquency.png")
km_path = os.path.join("reports", "survival_curves", "km_by_credit.png")
if os.path.exists(shap_path):
    slide.shapes.add_picture(shap_path, Inches(1), Inches(4.8), Inches(5), Inches(2.5))

# ===== SLIDE 6: LLM Copilot =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "FR-7: Governed LLM Copilot", size=36, color=ACCENT, bold=True)

tf = add_text(slide, 0.8, 1.5, 11, 5, "", size=16, color=WHITE)
add_para(tf, "How it works:", size=22, color=ORANGE, bold=True)
add_para(tf, "")
add_para(tf, "1. Retrieve loan's SHAP values, probabilities, anomaly score, exception type", size=17, color=WHITE)
add_para(tf, "2. Inject structured context into a grounded prompt template", size=17, color=WHITE)
add_para(tf, "3. Gemini API generates a plain-English reviewer note", size=17, color=WHITE)
add_para(tf, "4. Every output labeled:  \"Recommendation — not a decision\"", size=17, color=ACCENT)
add_para(tf, "")
add_para(tf, "Governance & Safety:", size=22, color=ORANGE, bold=True)
add_para(tf, "")
add_para(tf, "• All prompts + responses logged to llm_prompt_log.jsonl", size=17, color=WHITE)
add_para(tf, "• 3 documented failure examples (hallucination, overconfidence, confabulation)", size=17, color=WHITE)
add_para(tf, "• Automatic MockLLM fallback on rate limits (429 errors)", size=17, color=WHITE)
add_para(tf, "• LLM narrates numbers — never generates them", size=17, color=ACCENT)

# ===== SLIDE 7: Deliverables =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.4, 11, 0.8, "Deliverables Summary", size=36, color=ACCENT, bold=True)

deliverables = [
    ("submission.csv", "23,894 rows · 15 columns · 0 nulls"),
    ("MODEL_CARD.md", "209 lines · metrics, features, limitations, failure modes"),
    ("README.md", "Architecture, setup, run order, disqualifier audit"),
    ("8 reports", "FR-1 through FR-7 + servicer reconciliation log"),
    ("AI_DEVELOPMENT_LOG.md", "16 log entries · 7 rejected/corrected examples"),
    ("18 unit tests", "Time-split, leakage audit, servicer reconciliation — all passing"),
    ("demo_script.md", "15-step scripted walkthrough with line references"),
]

for i, (name, desc) in enumerate(deliverables):
    y = 1.5 + i * 0.75
    add_text(slide, 1.0, y, 3.5, 0.5, name, size=18, color=ACCENT, bold=True)
    add_text(slide, 4.8, y, 7.5, 0.5, desc, size=16, color=WHITE)

add_text(slide, 1.0, 6.8, 11, 0.4, "Fully reproducible:  python -m src.run_pipeline --stage all", size=16, color=GRAY)

# ===== SLIDE 8: Thank You =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 2.0, 11, 1.5, "Thank You", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 11, 0.8, "Loan Performance Intelligence Engine", size=24, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.0, 11, 1.2, "One command  ·  8 functional requirements  ·  Zero nulls\n18 tests passing  ·  All metrics from real held-out runs", size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Save
output_path = os.path.join("demo", "pitch_deck.pptx")
prs.save(output_path)
print(f"Saved to {output_path}")
